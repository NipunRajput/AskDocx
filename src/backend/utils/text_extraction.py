# utils/text_extraction.py
from pathlib import Path
import logging
import PyPDF2
import docx
import docx2txt
import openpyxl
from pptx import Presentation
from docx.opc.exceptions import PackageNotFoundError  # noqa

# --------------------------------------------------------------------- #
# Individual format helpers
# --------------------------------------------------------------------- #
def extract_text_from_pdf(path: str) -> str:
    try:
        reader = PyPDF2.PdfReader(path)
        return "\n".join((page.extract_text() or "") for page in reader.pages)
    except Exception as exc:
        logging.error("PDF read error: %s", exc)
        raise ValueError(f"PDF read error: {exc}")

def extract_text_from_docx(path: str) -> str:
    """Handle .docx first with python-docx, fall back to docx2txt."""
    try:
        return "\n".join(p.text for p in docx.Document(path).paragraphs)
    except Exception as err:
        logging.warning("python-docx failed (%s) – falling back to docx2txt", err)
        try:
            return docx2txt.process(path)
        except Exception as err2:
            logging.error("docx2txt also failed: %s", err2)
            raise ValueError(f"DOCX read error: {err}, then {err2}")

def extract_text_from_txt(path: str) -> str:
    for enc in ("utf-8", "latin-1", "cp1252"):
        try:
            with open(path, "r", encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    raise ValueError("TXT read error: unknown encoding")

def extract_text_from_xlsx(path: str) -> str:
    wb = openpyxl.load_workbook(path, data_only=True)
    rows = []
    for ws in wb.worksheets:
        for r in ws.iter_rows(values_only=True):
            rows.append("\t".join("" if c is None else str(c) for c in r))
    return "\n".join(rows)

def extract_text_from_pptx(path: str) -> str:
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

def extract_text_from_doc(path: str) -> str:
    """Old .doc (binary) – docx2txt can read it if `antiword` is installed."""
    try:
        return docx2txt.process(path)
    except Exception as exc:
        raise ValueError(f"DOC read error: {exc}")

# --------------------------------------------------------------------- #
# Unified dispatcher
# --------------------------------------------------------------------- #
def extract_text(path: str) -> str:
    ext = Path(path).suffix.lower()

    if ext == ".pdf":
        return extract_text_from_pdf(path)
    elif ext == ".docx":
        return extract_text_from_docx(path)
    elif ext == ".doc":
        return extract_text_from_doc(path)
    elif ext == ".txt":
        return extract_text_from_txt(path)
    elif ext == ".xlsx":
        return extract_text_from_xlsx(path)
    elif ext == ".pptx":
        return extract_text_from_pptx(path)

    raise ValueError(f"Unsupported file type: {ext}")
